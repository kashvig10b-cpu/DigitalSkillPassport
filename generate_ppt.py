import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def build_presentation(output_path):
    prs = Presentation()
    # 16:9 Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Modern Dark Theme Palette
    BG_DARK = RGBColor(11, 15, 25)         # Deep Slate 950
    CARD_BG = RGBColor(20, 27, 45)         # Slate 900 Card
    CARD_BORDER = RGBColor(38, 51, 80)     # Card Border
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)    # Slate 400
    TEXT_DIM = RGBColor(100, 116, 139)      # Slate 500
    PRIMARY = RGBColor(99, 102, 241)       # Indigo
    EMERALD = RGBColor(16, 185, 129)       # Emerald Green
    AMBER = RGBColor(245, 158, 11)         # Amber
    CYAN = RGBColor(6, 182, 212)           # Cyan Blue
    ROSE = RGBColor(244, 63, 94)           # Rose

    def set_slide_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="DIGITAL SKILL PASSPORT • FULL STACK"):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10.5)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(3)
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1)

    glow = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(4.9))
    glow.fill.solid()
    glow.fill.fore_color.rgb = CARD_BG
    glow.line.color.rgb = PRIMARY
    glow.line.width = Pt(2)

    tb = s1.shapes.add_textbox(Inches(1.4), Inches(1.7), Inches(10.5), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ENTERPRISE FULL-STACK ARCHITECTURE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Digital Skill Passport"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "Real-Time Verified Credential & Talent Discovery Platform with Anti-Scam Protection & QR Phone Verification"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(24)

    p = tf.add_paragraph()
    p.text = "Frontend: React 18, JavaScript (ES6+), Vite, Tailwind CSS • Backend: Node.js, Express, Socket.IO"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = "Database: MongoDB Atlas (Cloud) • Hosting: Vercel (Edge CDN) + Railway (Linux Microservice)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "Live Web Application: https://digital-skill-passport.vercel.app"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 2: THE PROBLEM
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "The Core Industry Problems We Solve", "PROBLEM STATEMENT")

    problems = [
        ("Resume Fraud & Fabrication", "Over 40% of resumes contain exaggerated or fake skills, certificates, and projects that traditional static PDF resumes cannot verify or disprove.", AMBER),
        ("Manual Verification Bottlenecks", "Universities and background-check companies spend weeks manually verifying diplomas, certificates, and test scores through slow email exchanges.", CYAN),
        ("Recruiter Scam Exploitation", "Unscreened third parties freely register on job boards as 'recruiters' to harvest student emails, phone numbers, and portfolios for spam and scams.", ROSE),
        ("Static, Unresponsive Records", "Paper and PDF resumes fail to reflect ongoing technical growth, verified credentials, or real-time skill progression.", PRIMARY)
    ]

    for i, (title, desc, color) in enumerate(problems):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.7 + row * 2.6)

        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s2.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), Inches(5.15), Inches(1.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"0{i+1}.  {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12.5)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: THREE-TIER FULL-STACK ARCHITECTURE
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3)
    add_header(s3, "Three-Tier Full-Stack Architecture Overview", "SYSTEM DESIGN")

    tiers = [
        ("TIER 1: CLIENT LAYER", "Frontend SPA (Vercel Edge)", "• Languages: JavaScript (ES6+), HTML5, CSS3\n• Framework: React.js 18 with Vite\n• Styling: Tailwind CSS & Lucide Icons\n• Dynamic Chart: Recharts (Radar Matrix)\n• Networking: Axios + Socket.IO Client\n• Global State: React Context API (Auth)", PRIMARY),
        ("TIER 2: API & REAL-TIME", "Backend Microservice (Railway)", "• Language: Node.js (v18+ Runtime)\n• Framework: Express.js (REST API)\n• WebSockets: Socket.IO (Event Bus)\n• Auth Security: JWT Tokens & Bcrypt\n• File Streaming: Multer Memory Storage\n• Architectural Pattern: MVC / Modular REST", CYAN),
        ("TIER 3: DATA & CLOUD", "Cloud Storage (MongoDB Atlas)", "• Database: MongoDB Atlas (M0 NoSQL)\n• ODM: Mongoose Schemas & Indexing\n• Data Models: 9 Collections (RBAC)\n• Persistent Storage: Binary BSON Buffer\n• Zero Disk Writes: Eliminates EACCES\n• Multi-University Data Isolation", EMERALD)
    ]

    for i, (tag, subtitle, desc, color) in enumerate(tiers):
        left = Inches(0.8 + i * 3.97)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3.8), Inches(5.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s3.shapes.add_textbox(left + Inches(0.3), Inches(2.0), Inches(3.2), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 4: FRONTEND ARCHITECTURE DEEP DIVE
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4)
    add_header(s4, "Frontend Architecture & Technologies", "CLIENT-SIDE STACK")

    fe_specs = [
        ("Core Languages", "JavaScript (ES6+ / JSX), HTML5, CSS3", "Modern JavaScript syntax, modular component architecture, and semantic HTML5.", PRIMARY),
        ("React 18 & Vite", "Single Page Application (SPA)", "Lightning-fast Hot Module Replacement (HMR) and optimized Rollup production bundling.", CYAN),
        ("Tailwind CSS", "Utility-First Responsive UI", "Custom dark-mode palette (`slate-950`), glowing border states, and mobile-friendly grids.", EMERALD),
        ("Recharts Matrix", "Multi-Axis Radar Skill Chart", "Dynamically renders candidate competencies across languages, frameworks, and skill levels.", AMBER),
        ("Socket.IO Client", "Bi-directional Real-Time Sync", "Listens for live certificate approvals, skill additions, and recruiter verifications.", PRIMARY),
        ("Axios HTTP Client", "Stateless REST API Engine", "Centralized API service with JWT bearer token interceptors and error handlers.", CYAN)
    ]

    for i, (title, sub, desc, color) in enumerate(fe_specs):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.7 + row * 1.7)

        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        tb = s4.shapes.add_textbox(left + Inches(0.25), top + Inches(0.15), Inches(5.25), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{title} • {sub}"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 5: BACKEND ARCHITECTURE DEEP DIVE
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5)
    add_header(s5, "Backend Architecture & Microservices", "SERVER-SIDE STACK")

    be_specs = [
        ("Runtime Environment", "Node.js (v18+ LTS)", "Asynchronous, non-blocking I/O event loop capable of thousands of concurrent requests.", PRIMARY),
        ("Web Framework", "Express.js REST API", "Modular route handlers, middleware pipelines, and strict request validation.", CYAN),
        ("Real-Time Event Gateway", "Socket.IO WebSockets Server", "Room-based event distribution: `admin_room`, `student_<id>`, and `passport_<id>` channels.", EMERALD),
        ("In-Memory File Streaming", "Multer (`memoryStorage`)", "Streams uploaded PDF/image buffers straight into MongoDB. 100% immune to container `EACCES` crashes.", AMBER),
        ("Authentication Engine", "JWT + Bcrypt.js Hashing", "Cryptographic JSON Web Tokens with 10-round salted password encryption.", ROSE),
        ("Security Middleware", "CORS & Rate Limiting", "Brute-force protection on auth endpoints and cross-origin resource isolation.", CYAN)
    ]

    for i, (title, sub, desc, color) in enumerate(be_specs):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.7 + row * 1.7)

        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        tb = s5.shapes.add_textbox(left + Inches(0.25), top + Inches(0.15), Inches(5.25), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{title} • {sub}"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 6: DATABASE & CLOUD STORAGE
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6)
    add_header(s6, "Database Modeling & Cloud Document Storage", "DATA TIER")

    # Left: Database Models
    c_left = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = CARD_BORDER

    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "MongoDB Atlas Schema Design"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_after = Pt(12)

    models = [
        ("User", "Auth, password hash, role (`student`, `recruiter`, `admin`), recruiterStatus"),
        ("StudentProfile", "Unique Passport ID, bio, academics, phone, profileCompletion score"),
        ("Skill", "Name, category, level, yearsOfExperience, studentId reference"),
        ("Project", "Title, techStack, description, githubUrl, liveUrl"),
        ("Certificate", "Title, issuer, issueDate, documentUrl, audit status (`VERIFIED`)"),
        ("Education & Experience", "Degrees, institutions, roles, durations, and achievements")
    ]
    for m_name, m_desc in models:
        p = tf.add_paragraph()
        p.text = f"• {m_name}: {m_desc}"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)

    # Right: Binary Cloud Storage
    c_right = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = CARD_BORDER

    tb2 = s6.shapes.add_textbox(Inches(7.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Binary Cloud File Storage (`FileAttachment`)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)

    p = tf2.add_paragraph()
    p.text = "• Zero Disk I/O (Railway Container Safe):\n  Resumes and certificate PDFs are streamed directly into MongoDB Atlas as BSON binary buffers (`Buffer`)."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(10)

    p = tf2.add_paragraph()
    p.text = "• Lifetime Document Persistence:\n  Files are never lost when server containers restart, update, or rebuild in the cloud."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(10)

    p = tf2.add_paragraph()
    p.text = "• High-Speed Streaming (`/api/files/:id`):\n  Native inline delivery with correct MIME types (`application/pdf`, `image/png`) for instant browser preview."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 7: SECURITY & ANTI-SCAM RECRUITER GATE
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7)
    add_header(s7, "Student Protection & Anti-Scam Recruiter Gate", "SECURITY ARCHITECTURE")

    steps_sec = [
        ("1. Company Verification on Registration", "When an employer registers, they must provide their verified Organization Name and official Website/LinkedIn link. Accounts start in `PENDING` status.", PRIMARY),
        ("2. HTTP 403 Forbidden Access Barrier", "The backend blocks unverified recruiters with a strict access barrier. Student emails, resumes, and portfolios are hidden to prevent data harvesting.", ROSE),
        ("3. 1-Click Verification Dashboard", "University Admins review employer credentials in the 'Recruiter Security' tab and can Approve or Revoke access with a single click.", AMBER),
        ("4. Instant WebSocket Real-Time Unlock", "The instant Admin approves an employer, Socket.IO triggers a live unlock on the recruiter's active tab without requiring a manual refresh.", EMERALD)
    ]

    for i, (title, desc, color) in enumerate(steps_sec):
        top = Inches(1.7 + i * 1.3)
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s7.shapes.add_textbox(Inches(1.1), top + Inches(0.12), Inches(11.1), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 8: PROFILE SCORING & MIN COMPLETION
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8)
    add_header(s8, "Dynamic Candidate Scoring & 'Min Completion' Filter", "TALENT DISCOVERY")

    c1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CARD_BORDER

    tb = s8.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Weighted Scoring Formula (100% Total)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_after = Pt(12)

    weights = [
        ("Technical Skills Matrix Added", "+15%"),
        ("Live GitHub Projects Attached", "+15%"),
        ("Education Academic History", "+15%"),
        ("Degree & Department Info", "+15%"),
        ("Uploaded Resume or Certificate", "+10%"),
        ("Verified Phone & Location", "+10%"),
        ("Detailed Professional Bio", "+10%"),
        ("Student Profile Photo", "+10%")
    ]
    for w_name, w_val in weights:
        p = tf.add_paragraph()
        p.text = f"• {w_name.ljust(35)} {w_val}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(3)

    c2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = CARD_BORDER

    tb2 = s8.shapes.add_textbox(Inches(7.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "The 'Min Completion' Recruiter Slider"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(12)

    p = tf2.add_paragraph()
    p.text = "• Eliminates Incomplete Accounts:\n  Recruiters filter out abandoned accounts by dragging the slider (e.g. 70% or 80%)."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(10)

    p = tf2.add_paragraph()
    p.text = "• Dynamic Candidate Matching:\n  Scores are calculated on-the-fly and synced to MongoDB so candidates reflect their true progress."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(10)

    p = tf2.add_paragraph()
    p.text = "• Verified Only Filter:\n  A 1-click checkbox to only display candidates whose certificates have been audited by the University Admin."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 9: DEVOPS & CI/CD PIPELINE
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9)
    add_header(s9, "DevOps, Cloud Hosting & Automated CI/CD Pipeline", "DEPLOYMENT WORKFLOW")

    devops_cards = [
        ("Git Version Control", "GitHub (`main` branch)", "• Centralized source of truth\n• Atomic commits with verified code linting\n• Webhook integration for instant deployment", PRIMARY),
        ("Vercel Edge Deployment", "Frontend Hosting", "• Global CDN caching with sub-100ms latency\n• Automated Vite build on push (~20s)\n• Single Page App rewrites (`vercel.json`)", CYAN),
        ("Railway Microservice", "Backend Cloud Hosting", "• Containerized Linux environment\n• Zero-downtime rolling deploys (~35s)\n• Automatic crash restart & cloud health check", EMERALD)
    ]

    for i, (title, sub, desc, color) in enumerate(devops_cards):
        left = Inches(0.8 + i * 3.97)
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(3.8), Inches(5.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s9.shapes.add_textbox(left + Inches(0.3), Inches(2.0), Inches(3.2), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(3)

        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 10: END-TO-END DEMO WORKFLOW
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10)
    add_header(s10, "Step-by-Step Live Demonstration Workflow", "PLATFORM DEMO")

    workflow_steps = [
        ("Step 1: Student Onboarding", "Student registers, adds technical skills, links GitHub projects, and uploads their official resume.", PRIMARY),
        ("Step 2: Admin Audit Queue", "Admin inspects submitted certificate PDF proofs, verifying authenticity with 1-click approval.", AMBER),
        ("Step 3: Recruiter Discovery", "Recruiter filters candidates by skill, sets Min Completion to 60%, and connects via 1-click Gmail.", CYAN),
        ("Step 4: Phone Camera QR Scan", "Anyone scans the student's QR code on mobile to load the live, verified digital skill passport anywhere in the world.", EMERALD)
    ]

    for i, (title, desc, color) in enumerate(workflow_steps):
        left = Inches(0.8 + i * 2.97)
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.7), Inches(2.8), Inches(5.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s10.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(2.4), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 11: IMPACT & FUTURE ROADMAP
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11)
    add_header(s11, "Project Impact & Future Enhancements", "CONCLUSION")

    c1 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = CARD_BORDER

    tb = s11.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Project Impacts"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "• Trust & Authenticity:\n  Eliminates credential fraud via official university administrative stamps.\n\n• Zero-Friction Discovery:\n  Instant QR code scanning opens verified credentials on any phone globally.\n\n• Student Privacy Protection:\n  Administrative gate blocks scam employers from accessing student contacts."
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MUTED

    c2 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = CARD_BORDER

    tb2 = s11.shapes.add_textbox(Inches(7.1), Inches(1.95), Inches(5.1), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Future Technical Scope"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_after = Pt(12)

    p = tf2.add_paragraph()
    p.text = "• Blockchain Soulbound Credentials:\n  Minting verified degrees and certificates as non-transferable NFTs on Polygon.\n\n• Automated Code Skill Assessments:\n  Integrated coding challenges that automatically accredit skills upon passing.\n\n• University ERP System Webhooks:\n  Direct API synchronization with student information systems (SIS)."
    p.font.size = Pt(12.5)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 12: THANK YOU & Q&A
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12)

    box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = EMERALD
    box.line.width = Pt(2)

    tb = s12.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.text = "Questions & Discussion"
    p.font.size = Pt(22)
    p.font.color.rgb = CYAN
    p.space_after = Pt(24)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.text = "Live Deployment: https://digital-skill-passport.vercel.app\nGitHub Repository: github.com/kashvig10b-cpu/DigitalSkillPassport"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED

    prs.save(output_path)
    print(f"Updated presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "Digital_Skill_Passport_Presentation.pptx"
    build_presentation(out)
